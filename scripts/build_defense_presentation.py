#!/usr/bin/env python3
"""Defense presentation: plain layout (text left / image right) + algorithm block. 16:9."""

from __future__ import annotations

import argparse
import sys
import tempfile
from dataclasses import dataclass
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    import subprocess

    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

try:
    from PIL import Image
except ImportError:
    Image = None  # type: ignore

REPO = Path(__file__).resolve().parents[1]
OUT = REPO / "docs" / "presentation" / "defense_presentation.pptx"

W = Inches(13.333)
H = Inches(7.5)

# Plain style — без декора
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0x88, 0x88, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = BLACK
MUTED = GRAY
ACCENT = BLACK

# Layout (inches)
MX = Inches(0.72)
MY_TITLE = Inches(0.52)
MY_BODY = Inches(1.22)
CONTENT_W = Inches(11.893)
HALF_W = Inches(6.15)
COL_GAP = Inches(0.28)
RIGHT_X = MX + HALF_W + COL_GAP

ASSETS = REPO / "docs" / "diploma_drafts" / "assets"
LAB = REPO / "results" / "lab"
RPT = LAB / "presentation_report"
FLC = LAB / "final_large_complex"
IDEAL = LAB / "ideal_coverage_demo"
DYN = LAB / "presentation_dynamic"
DYN_ML = LAB / "presentation_dynamic_ml_rl"
STAT = LAB / "presentation_static"

_image_cache: dict[tuple, Path] = {}
_temp_files: list[Path] = []


@dataclass(frozen=True)
class FontProfile:
    main: int
    caption: int
    slide_title: int
    title_main: int
    title_sub: int
    title_info: int
    appendix_title: int
    section_title: int
    table_header: int
    table_body: int
    figure_label: int
    thanks_head: int
    thanks_body: int
    stat_big: int
    stat_label: int


FONT_DEFAULT = FontProfile(16, 11, 22, 28, 17, 14, 15, 24, 11, 10, 11, 34, 15, 36, 12)
FONT_MAIN14 = FontProfile(15, 11, 20, 26, 15, 13, 14, 18, 11, 11, 10, 32, 14, 28, 11)

FONT = FONT_MAIN14

# Подписи на картах покрытия / GIF (крупнее и жирнее)
FIG_LABEL_SIZE = 15
FIG_LABEL_BOLD = True

# Справочник алгоритмов для приложения Б
ALGO_CATALOG: list[tuple[str, str]] = [
    ("baseline_random_walk", "случайный обход; минимальная база для сравнения"),
    ("baseline_grid", "последовательный обход сетки; детерминированно, без координации"),
    ("baseline_voronoi", "разбиение площади на зоны роботов; склад, патруль"),
    ("baseline_frontier", "движение к границе непокрытого; быстрый рост покрытия"),
    ("boustrophedon", "«газонокосилка»; системные полосы на известной карте"),
    ("stc", "обход по остовному дереву; полнота на static-сетке"),
    ("darp_boustro", "баланс длин путей между роботами (DARP + boustrophedon)"),
    ("darp_stc", "баланс нагрузки + полный обход STC"),
    ("cbs_boustrophedon", "CBS-де-конфликт поверх boustrophedon"),
    ("cbs_darp_boustro", "CBS-де-конфликт поверх DARP"),
    ("ppo_policy", "политика обучения с подкреплением (PPO)"),
    ("ml_guided", "нейросеть выбирает шаг; режимы с ограничителем безопасности"),
    ("ml_goal", "нейросеть выбирает цель покрытия; несколько режимов"),
]


def _fig_label(slide, left, top, width, text: str):
    _text(
        slide, left, top, width, Inches(0.32), text,
        size=FIG_LABEL_SIZE, bold=FIG_LABEL_BOLD, color=BLACK, align=PP_ALIGN.CENTER,
    )


def _exists(p: Path | None) -> Path | None:
    return p if p and p.is_file() else None


def _pick(*candidates: Path) -> Path | None:
    for c in candidates:
        if p := _exists(c):
            return p
    return None


def _flc(algo: str, kind: str = "coverage_2d", seed: int = 0) -> Path | None:
    return _exists(FLC / f"large_complex_dynamic__{algo}__seed{seed}_{kind}.png")


def _ideal(algo: str, kind: str = "coverage_2d", seed: int = 0) -> Path | None:
    return _exists(IDEAL / f"large_complex_dynamic__{algo}__seed{seed}_{kind}.png")


def _dyn_gif(algo: str, seed: int = 0) -> Path | None:
    return _exists(DYN / f"dynamic_B_long__{algo}__seed{seed}.gif")


def _static_gif(algo: str, seed: int = 0) -> Path | None:
    return _exists(STAT / f"static_A_long__{algo}__seed{seed}.gif")


def _static_2d(algo: str, seed: int = 0) -> Path | None:
    return _exists(STAT / f"static_A_long__{algo}__seed{seed}_coverage_2d.png")


def _dyn_png(algo: str, kind: str = "coverage_2d", seed: int = 0) -> Path | None:
    return _exists(DYN / f"dynamic_B_long__{algo}__seed{seed}_{kind}.png")


def _ideal_gif(algo: str, seed: int = 0) -> Path | None:
    return _exists(IDEAL / f"large_complex_dynamic__{algo}__seed{seed}.gif")


def _dyn_ml_gif(seed: int = 0) -> Path | None:
    return _exists(DYN_ML / f"dynamic_B_long__ml_guided__seed{seed}.gif")


def _dyn_ml_png(seed: int = 0) -> Path | None:
    return _exists(DYN_ML / f"dynamic_B_long__ml_guided__seed{seed}_coverage_2d.png")


def _rel(p: Path | None) -> str:
    if not p:
        return ""
    try:
        return str(p.relative_to(REPO))
    except ValueError:
        return str(p)


def _gif_to_png(gif: Path) -> Path | None:
    if not Image:
        return None
    key = ("gif_frame", str(gif))
    if key in _image_cache:
        return _image_cache[key]
    try:
        im = Image.open(gif)
        im.seek(0)
        if im.mode not in ("RGB", "L"):
            im = im.convert("RGB")
        fd, name = tempfile.mkstemp(suffix=".png", prefix="defense_")
        out = Path(name)
        im.save(out, "PNG", optimize=True)
        _temp_files.append(out)
        _image_cache[key] = out
        return out
    except OSError:
        return None


def _prepare_media(path: Path | None, *, max_px: int = 1400, allow_gif: bool = False) -> Path | None:
    if not path or not path.is_file():
        return None
    cache_key = (str(path), max_px, allow_gif)
    if cache_key in _image_cache:
        return _image_cache[cache_key]

    src = path
    if path.suffix.lower() == ".gif":
        if allow_gif:
            _image_cache[cache_key] = path
            return path
        stem = path.stem
        for alt in (
            path.parent / f"{stem}_coverage_2d.png",
            path.parent / f"{stem}.png",
        ):
            if alt.is_file():
                src = alt
                break
        else:
            src = _gif_to_png(path)
            if not src:
                return None

    if not Image or src.suffix.lower() not in (".png", ".jpg", ".jpeg"):
        _image_cache[cache_key] = src
        return src

    try:
        im = Image.open(src)
        im.load()
        w, h = im.size
        if max(w, h) > max_px:
            ratio = max_px / max(w, h)
            im = im.resize((int(w * ratio), int(h * ratio)), Image.Resampling.LANCZOS)
        if im.mode in ("RGBA", "P"):
            bg = Image.new("RGB", im.size, (255, 255, 255))
            if im.mode == "P":
                im = im.convert("RGBA")
            bg.paste(im, mask=im.split()[-1] if im.mode == "RGBA" else None)
            im = bg
        elif im.mode != "RGB":
            im = im.convert("RGB")
        fd, name = tempfile.mkstemp(suffix=".jpg", prefix="defense_")
        out = Path(name)
        im.save(out, "JPEG", quality=82, optimize=True)
        _temp_files.append(out)
        _image_cache[cache_key] = out
        return out
    except OSError:
        _image_cache[cache_key] = src
        return src


A = {
    "panel_static": _pick(RPT / "panel_median_static.png"),
    "panel_dynamic": _pick(RPT / "panel_median_dynamic.png"),
    "panel_seed0": _pick(RPT / "panel_seed0_static_dynamic.png"),
    "panel_flc": _pick(FLC / "panel_final_large_complex.png", IDEAL / "panel_final_large_complex.png"),
    "flc_frontier_2d": _flc("baseline_frontier"),
    "flc_voronoi_2d": _flc("baseline_voronoi"),
    "flc_darp_2d": _flc("darp_boustro"),
    "flc_stc_2d": _flc("stc"),
    "flc_ml_2d": _flc("ml_guided"),
    "flc_frontier_time": _flc("baseline_frontier", "coverage_vs_time"),
    "flc_darp_time": _flc("darp_boustro", "coverage_vs_time"),
    "ideal_guarded_2d": _ideal("ml_guided_guarded", seed=2),
    "ideal_guarded_gif": _ideal_gif("ml_guided_guarded", seed=2),
    "isaac": _pick(
        IDEAL / "large_complex_dynamic__ml_guided_guarded__seed2_isaac.gif",
        IDEAL / "_smoke_seed2_isaac.gif",
    ),
    "dyn_frontier_gif": _dyn_gif("baseline_frontier"),
    "dyn_voronoi_gif": _dyn_gif("baseline_voronoi"),
    "dyn_darp_gif": _dyn_gif("darp_boustro"),
    "dyn_stc_gif": _dyn_gif("stc"),
    "dyn_ml_gif": _dyn_ml_gif(0),
    "dyn_frontier_png": _dyn_png("baseline_frontier"),
    "dyn_ml_png": _dyn_ml_png(0),
    "ideal_frontier_gif": _ideal_gif("baseline_frontier", 0),
    "ideal_voronoi_gif": _ideal_gif("baseline_voronoi", 0),
    "ideal_darp_gif": _ideal_gif("darp_boustro", 0),
    "ideal_stc_gif": _ideal_gif("stc", 0),
    "stat_frontier_gif": _static_gif("baseline_frontier"),
    "stat_voronoi_gif": _static_gif("baseline_voronoi"),
    "stat_frontier_2d": _static_2d("baseline_frontier"),
    "stat_voronoi_2d": _static_2d("baseline_voronoi"),
    "report_cards": _pick(RPT / "report_cards.md"),
    "sch_frontier": _pick(ASSETS / "algo_schematic_frontier.png"),
    "sch_voronoi": _pick(ASSETS / "algo_schematic_voronoi.png"),
    "sch_stc": _pick(ASSETS / "algo_schematic_stc.png"),
    "sch_darp": _pick(ASSETS / "algo_schematic_darp.png"),
    "sch_cbs": _pick(ASSETS / "algo_schematic_cbs.png"),
    "sch_ml": _pick(ASSETS / "algo_schematic_ml_guarded.png"),
    "sch_classification": _pick(ASSETS / "classification_layers.png"),
    "sch_algos": _pick(ASSETS / "classification_algorithms.png"),
    "frame_guarded": _pick(ASSETS / "frame_ml_guided_guarded_seed2.png"),
    "frame_isaac": _pick(ASSETS / "frame_isaac_replay_seed2.png"),
    "openvla": _pick(ASSETS / "openvla_architecture.jpg"),
}


def _viz(embed_gif: bool, gif: Path | None, *fallbacks: Path | None) -> Path | None:
    """При --embed-gif предпочитаем анимацию, иначе PNG/2D."""
    if embed_gif and gif:
        return gif
    for f in fallbacks:
        if f:
            return f
    return gif


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, left, top, width, height, fill: RGBColor, border: RGBColor | None = None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    return sh


def _code_footer(slide, code: str):
    if code:
        _text(slide, Inches(12.35), Inches(7.06), Inches(0.85), Inches(0.25), code, size=9, color=LIGHT, align=PP_ALIGN.RIGHT)


def _text(slide, left, top, width, height, text, size=None, bold=False, color=INK, align=PP_ALIGN.LEFT):
    size = size or FONT.main
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align


def _bullets(slide, left, top, width, height, lines: list[str], size=None):
    size = size or FONT.main
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines[:9]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line if line.startswith("·") else f"·  {line}"
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p.space_after = Pt(10)
        p.line_spacing = 1.15


def _caption(slide, text: str):
    if text:
        _text(slide, MX, Inches(7.0), CONTENT_W, Inches(0.3), text, size=FONT.caption, color=LIGHT)


def _notes(slide, text: str):
    if text.strip():
        slide.notes_slide.notes_text_frame.text = text.strip()


def _img(slide, path: Path | None, left, top, width, height, *, allow_gif: bool = False) -> bool:
    prepared = _prepare_media(path, allow_gif=allow_gif)
    if prepared and prepared.is_file():
        slide.shapes.add_picture(str(prepared), left, top, width=width, height=height)
        return True
    return False


def _heading(slide, text: str):
    _text(slide, MX, MY_TITLE, CONTENT_W, Inches(0.55), text, size=FONT.slide_title, bold=True, color=BLACK)


def _slide_text_image(
    prs: Presentation,
    title: str,
    lines: list[str],
    image: Path | None,
    *,
    allow_gif: bool = False,
    footnote: str = "",
    notes: str = "",
    bullet_size: int | None = None,
):
    """Обычный слайд: текст слева, одна картинка справа."""
    s = _blank(prs)
    _heading(s, title)
    _bullets(s, MX, MY_BODY, HALF_W, Inches(5.8), lines, size=bullet_size or FONT.main)
    _img(s, image, RIGHT_X, MY_BODY, HALF_W, Inches(5.55), allow_gif=allow_gif)
    if footnote:
        _caption(s, footnote)
    if notes:
        _notes(s, notes)
    return s


def _slide_image_main(
    prs: Presentation,
    title: str,
    image: Path | None,
    *,
    caption: str = "",
    notes: str = "",
):
    """Слайд с крупной схемой на всю ширину."""
    s = _blank(prs)
    _heading(s, title)
    if not _img(s, image, MX, MY_BODY, CONTENT_W, Inches(5.85)):
        _text(s, MX, Inches(3.2), CONTENT_W, Inches(0.5), "Схема не найдена", color=MUTED)
    if caption:
        _caption(s, caption)
    if notes:
        _notes(s, notes)
    return s


def _slide_table(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    *,
    caption: str = "",
    notes: str = "",
    col_widths: list[float] | None = None,
):
    """Инженерный слайд: заголовок + широкая таблица с числами."""
    s = _blank(prs)
    _heading(s, title)
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_h = Inches(0.42 * n_rows + 0.1)
    tbl = s.shapes.add_table(n_rows, n_cols, MX, MY_BODY, CONTENT_W, table_h).table
    if col_widths:
        total = sum(col_widths)
        for j, w in enumerate(col_widths):
            tbl.columns[j].width = Inches(11.893 * w / total)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        for p in c.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(14)
                r.font.bold = True
                r.font.color.rgb = BLACK
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = str(val)
            for p in c.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.size = Pt(13)
                    r.font.color.rgb = INK
                    if j == 0:
                        r.font.bold = True
    if caption:
        _caption(s, caption)
    if notes:
        _notes(s, notes)
    return s


def _slide_gif_grid(
    prs: Presentation,
    title: str,
    items: list[tuple[str, Path | None]],
    *,
    cols: int = 3,
    code: str = "",
    notes: str = "",
):
    """Сетка GIF по сцене: подпись под каждым."""
    s = _blank(prs)
    _heading(s, title)
    if not items:
        return s
    gap = Inches(0.18)
    n = len(items)
    rows = (n + cols - 1) // cols
    cw = Inches(3.85)
    ch = Inches(2.45) if rows > 1 else Inches(5.2)
    for i, (label, path) in enumerate(items):
        c, r = i % cols, i // cols
        l = MX + c * (cw + gap)
        t = MY_BODY + r * (ch + gap)
        _img(s, path, l, t, cw, ch - Inches(0.28), allow_gif=True)
        _fig_label(s, l, t + ch - Inches(0.30), cw, label)
    if code:
        _code_footer(s, code)
    if notes:
        _notes(s, notes)
    return s


def _app_columns(
    prs: Presentation,
    code: str,
    title: str,
    columns: list[tuple[str, list[str]]],
):
    """Подвал: несколько разделов в столбцах (A | B | …)."""
    s = _blank(prs)
    _heading(s, title)
    n = len(columns)
    gap = Inches(0.35)
    cw = Inches((12.0 - 0.72 * 2 - (n - 1) * 0.35) / n)  # noqa: rough — use Inches math
    cw = Inches(3.75) if n >= 3 else HALF_W
    for i, (head, lines) in enumerate(columns):
        x = MX + i * (cw + gap)
        _text(s, x, MY_BODY, cw, Inches(0.32), head, size=FONT.main, bold=True)
        _bullets(s, x, MY_BODY + Inches(0.38), cw, Inches(5.3), lines, size=11)
    _code_footer(s, code)


def _slide_algo_promo(prs: Presentation):
    """Промо: все схемы планировщиков из assets."""
    s = _blank(prs)
    _heading(s, "Планировщики покрытия в конструкторе")
    items = [
        ("Frontier", A["sch_frontier"]),
        ("Voronoi · разбиение", A["sch_voronoi"]),
        ("STC", A["sch_stc"]),
        ("DARP", A["sch_darp"]),
        ("CBS · де-конфликт", A["sch_cbs"]),
        ("ML + guardrail", A["sch_ml"]),
    ]
    cols, rows = 3, 2
    cw = Inches(3.95)
    ch = Inches(2.55)
    g = Inches(0.22)
    for i, (label, path) in enumerate(items):
        c, r = i % cols, i // cols
        l = MX + c * (cw + g)
        t = MY_BODY + r * (ch + g)
        _img(s, path, l, t, cw, ch - Inches(0.28))
        _fig_label(s, l, t + ch - Inches(0.30), cw, label)
    _caption(s, "docs/diploma_drafts/assets/algo_schematic_*.png")
    _notes(s, "25–30 с. Шесть модулей покрытия — дальше «для чего» на след. слайдах.")


def _slide_algo_detail(
    prs: Presentation,
    title: str,
    left_name: str,
    left_sch: Path | None,
    left_cov: Path | None,
    left_lines: list[str],
    right_name: str,
    right_sch: Path | None,
    right_cov: Path | None,
    right_lines: list[str],
    note: str = "",
):
    s = _blank(prs)
    _heading(s, title)
    sch_h = Inches(2.45)
    cov_h = Inches(2.55)
    cov_t = MY_BODY + sch_h + Inches(0.12)

    def _col(name, sch, cov, lines, left):
        x = MX if left else RIGHT_X
        _fig_label(s, x, MY_BODY - Inches(0.04), HALF_W, name)
        _img(s, sch, x, MY_BODY + Inches(0.22), HALF_W, sch_h)
        bullet_t = cov_t
        if cov:
            _img(s, cov, x, cov_t, HALF_W, cov_h)
            bullet_t = cov_t + cov_h + Inches(0.08)
        else:
            bullet_t = MY_BODY + sch_h + Inches(0.35)
        _bullets(s, x, bullet_t, HALF_W, Inches(1.85), lines, size=FONT.main - 1)

    _col(left_name, left_sch, left_cov, left_lines, True)
    _col(right_name, right_sch, right_cov, right_lines, False)
    if note:
        _caption(s, note)
    _notes(s, f"{title} · по ~20 с на колонку.")


# --- Main deck ---


def build_main(prs: Presentation, *, embed_gif: bool):
    half = HALF_W
    gap = COL_GAP
    right_x = RIGHT_X

    # 1 Title
    s = _blank(prs)
    _text(s, MX, Inches(2.0), Inches(11.5), Inches(1.4),
          "Многоагентная навигация и координация наземных роботов",
          size=FONT.title_main, bold=True)
    _text(s, MX, Inches(3.5), Inches(11.5), Inches(0.45),
          "Платформа сравнения multi-robot coverage",
          size=FONT.title_sub)
    _text(s, MX, Inches(5.2), Inches(11.5), Inches(0.8),
          "Кирюшин Артём · М3О-412Б-22 · каф. 307 · МАИ 2026\nРуководитель: Склеймин Ю.Б.",
          size=FONT.title_info, color=GRAY)
    _notes(s, "Титул 15 с.")

    # 2 Цель и задачи — общими буллетами
    _slide_text_image(
        prs, "Цель и задачи",
        [
            "Цель: многоагентная система навигации и координации наземных роботов",
            "с воспроизводимым сравнением алгоритмов покрытия",
            "Задачи:",
            "  выбрать и обосновать симуляционную среду",
            "  проанализировать подходы к навигации и координации",
            "  спроектировать архитектуру системы (слои, модули, метрики)",
            "  реализовать и интегрировать систему в симуляторе",
            "  провести тесты и сравнить результаты на единых сценах",
            "  сформулировать итоги и рекомендации по выбору конфигурации",
        ],
        A["sch_classification"],
        bullet_size=FONT.main - 1,
        notes="45 с. Задачи — общими формулировками, без таблицы плана.",
    )

    # 3 Актуальность · 2025 → 2026
    _slide_text_image(
        prs, "Актуальность · рынок 2025–2026",
        [
            "AMR (автономные мобильные роботы):",
            "  2025 — ~$5,8 млрд  →  2026 — ~$6,8 млрд (+18% г/г)",
            "Роботы на складах и в логистике — основной драйвер роста",
            "Автоматизация складов (гл.): ~$29 млрд (2025), к 2026 ~$33 млрд",
            "Сервисная робототехника: ~$28 млрд (2025), рост ~15% в год",
            "К 2026: >40% новых складов планируют группы AMR (оценки аналитиков)",
            "В РФ: пилоты беспилотной логистики и патрулирования на объектах",
            "Проблема: в статьях нет единого стенда — нельзя честно сравнить",
        ],
        _viz(embed_gif, A["ideal_guarded_gif"], A["ideal_guarded_2d"]),
        allow_gif=embed_gif,
        footnote="* The Business Research Company (2026); Fact.MR; MarketsandMarkets",
        bullet_size=FONT.main - 1,
        notes="50 с. Цифры за 2025, прогноз на 2026; место под график.",
    )

    # 4 Novelty — авторская идея
    _slide_text_image(
        prs, "Новизна работы",
        [
            "Вклад: прикладная классификация, не «ещё один planner»",
            "5 осей — по каким вопросам выбирать метод",
            "5 слоёв — что комбинировать в YAML",
            "Согласовано с кодом (_algo_factory) и batch",
            "Стенд: честное сравнение в одном контуре",
        ],
        A["panel_flc"],
        notes="50 с. Ключ: я предложил рамку выбора под задачу.",
    )

    # 5 Classification — 5 слоёв (нейтральная схема)
    _slide_image_main(
        prs,
        "Классификация · 5 слоёв конструктора",
        A["sch_classification"],
        caption="",
        notes="45 с. Главная схема: пять слоёв, нейтрально.",
    )

    # 6 Три семейства алгоритмов (метрики — слой 5, не семейство)
    _slide_image_main(
        prs,
        "Три семейства алгоритмов",
        A["sch_algos"],
        caption="метрики и ограничители — в слое «Проверка», это не отдельное семейство",
        notes="40 с. Три группы: классика, координация, обучение.",
    )

    # 7 Промо планировщиков (6 схем)
    _slide_algo_promo(prs)

    # 8 Платформа · цифры
    _slide_text_image(
        prs, "Платформа multi_robot_coverage · цифры",
        [
            "15+ политик: 4 baseline · STC · DARP × 2 · CBS-обёртка · PPO",
            "ML: ml_guided × 3 режима · ml_goal × 6 режимов",
            "3 семьи сцен: static_A / dynamic_B / large_complex",
            "Основная серия batch (один прогон = политика + сцена + повтор):",
            "  static: 5 политик × 5 повторов = 25",
            "  dynamic: 5 политик × 5 повторов = 25",
            "  stress: 5 политик × 3 повтора = 15",
            "  итого: 65 прогонов (25 + 25 + 15)",
            "8 метрик: покрытие · время · путь · блокировки · пешеходы · нагрузка",
            "Контракт JSON: coverage_lab → run_batch → summary.csv",
        ],
        A["panel_dynamic"],
        bullet_size=FONT.main - 1,
        notes="50 с. 65 = 25+25+15 (static+dynamic+large_complex). Не говорить 150.",
    )

    # 9 Методология
    _slide_text_image(
        prs, "Методология сравнения",
        [
            "static_A_long — целевое 0.90, 4 робота",
            "dynamic_B_long — целевое 0.80, 6 пешеходов",
            "large_complex — 8 роботов, узкие проходы",
            "Семя 0, 1, 2 — медиана по seed",
            "Каждый прогон даёт coverage_2d, time, JSON",
        ],
        A["panel_static"],
        notes="35 с.",
    )

    # 10 Результаты · инженерная таблица
    _slide_table(
        prs, "Результаты large_complex_dynamic · медиана seed",
        ["алгоритм", "покр.", "время до цели, с", "путь, м", "блок.", "наруш. пед.", "равном. нагр."],
        [
            ["baseline_frontier", "0.78", "87.9",  "1393", "1459",  "26",  "0.34"],
            ["baseline_voronoi",  "0.78", "79.2",  "1336", "1102",  "32",  "0.33"],
            ["stc",               "0.78", "255.7", "1990", "12707", "205", "0.50"],
            ["darp_boustro",      "0.69", "—",     "10357","172",   "204", "0.01"],
            ["ml_guided",         "0.78", "79.2",  "1336", "1102",  "32",  "0.33"],
        ],
        col_widths=[2.0, 0.7, 1.0, 1.0, 0.8, 1.1, 1.2],
        caption="равном. нагр. — коэф. вариации длин путей роботов (меньше = равнее нагрузка)",
        notes="60 с. Главный инженерный слайд — на нём числа и читаем.",
    )

    # 11 Применимость · одна страница с выводами по числам
    _slide_table(
        prs, "Применимость алгоритмов · выбор по сценарию",
        ["алгоритм", "сильная сторона", "слабая сторона", "когда брать"],
        [
            ["baseline_frontier",   "быстрый рост покрытия (88 с)", "средние наруш. у пешеходов", "разведка, частичная карта"],
            ["baseline_voronoi",    "зоны, меньше нарушений",       "слабее в узких проходах",    "склад, патруль с людьми"],
            ["stc",                 "полнота на сетке",             "медленно: 256 с до цели, 12 700 блокировок", "известная статич. карта"],
            ["darp_boustro",        "равная нагрузка (коэф. 0.01)", "путь 10 357 м",              "много роботов, равномерно"],
            ["ml_guided_guarded",   "покрытие + ограничитель",      "≈ voronoi без обучения",     "нужны покрытие и безопасность"],
            ["cbs_* / ppo_policy",  "снятие столкновений / RL",     "надстройки",                 "плотный трафик / динамика"],
        ],
        col_widths=[1.7, 2.4, 2.2, 2.6],
        notes="50 с. Это слайд, на котором повторяем «выбор под задачу».",
    )

    # 12 GIF — main visualization: large_complex + Isaac
    if embed_gif:
        _slide_gif_grid(
            prs,
            "Стресс-сцена large_complex_dynamic · Isaac Sim",
            [
                ("baseline_frontier", A["ideal_frontier_gif"]),
                ("baseline_voronoi", A["ideal_voronoi_gif"]),
                ("darp_boustro", A["ideal_darp_gif"]),
                ("stc", A["ideal_stc_gif"]),
                ("ml_guided_guarded", A["ideal_guarded_gif"]),
                ("Isaac Sim", A["isaac"]),
            ],
            cols=3,
            notes="40 с. Один слайд GIF в основной части. Static/dynamic — в подвал.",
        )
    else:
        _slide_text_image(
            prs, "Стресс-сцена (PNG)",
            ["GIF: python scripts/build_defense_presentation.py --embed-gif"],
            A["dyn_frontier_png"],
            notes="30 с.",
        )

    # 16 VLA
    _slide_text_image(
        prs, "Перспектива · VLA",
        [
            "Верхний уровень: выбор модуля конструктора",
            "Нижний уровень: coverage_lab + ограничитель безопасности",
            "Обзор, не эксперимент ВКР",
        ],
        A["openvla"],
        notes="25 с.",
    )

    # 14 Итоги и результаты
    _slide_text_image(
        prs, "Итоги и результаты",
        [
            "Сделано:",
            "  Классификация (5 слоёв) и сравнение 15+ политик на batch",
            "  Платформа: coverage_lab, run_batch, JSON, отчёты",
            "  Репозиторий с кодом, сценами YAML, results/lab/",
            "  Панели, GIF, Isaac Sim (3D), схемы алгоритмов",
            "Результаты:",
            "  Паттерны по метрикам: voronoi 79 с vs STC 256 с при том же %",
            "  DARP — наиболее равная нагрузка на роботов (коэф. 0.01)",
            "  Один % покрытия не заменяет панель KPI",
            "Дальше: Isaac, VLA, ML-goal, новые сцены в том же git",
        ],
        A["panel_flc"],
        bullet_size=FONT.main - 2,
        notes="50 с. Итоги + что в репозитории.",
    )

    # 15 Thanks
    s = _blank(prs)
    _text(s, MX, Inches(2.8), Inches(11.5), Inches(0.7), "Спасибо за внимание",
          size=FONT.thanks_head, bold=True, align=PP_ALIGN.CENTER)
    _text(s, MX, Inches(3.7), Inches(11.5), Inches(2.5),
          "Вопросы — приложение (после слайда 15)\n\n"
          "А схемы  ·  Б алгоритмы  ·  В сцены  ·  Г карты  ·  Д цифры",
          size=FONT.thanks_body, color=GRAY, align=PP_ALIGN.CENTER)
    _notes(s, "Указать подвал.")


# --- Compact appendix ---


def _parse_report_cards() -> list[list[str]]:
    path = A["report_cards"]
    if not path:
        return []
    rows = []
    for line in path.read_text(encoding="utf-8").splitlines():
        if not line.startswith("|") or "---" in line or "algorithm" in line:
            continue
        cells = [c.strip().strip("`") for c in line.split("|")[1:-1]]
        if cells:
            rows.append(cells)
    return rows[:8]


def _app_compact_section(prs: Presentation, code: str, title: str, body: list[str]):
    s = _blank(prs)
    _heading(s, title)
    _bullets(s, MX, MY_BODY, CONTENT_W, Inches(5.8), body)
    _code_footer(s, code)


def _app_algo_catalog(prs: Presentation, code: str, title: str):
    """Приложение Б: имя алгоритма + краткое описание (2 колонки)."""
    s = _blank(prs)
    _heading(s, title)
    mid = len(ALGO_CATALOG) // 2 + len(ALGO_CATALOG) % 2
    gap = Inches(0.35)
    for col_i, chunk in enumerate([ALGO_CATALOG[:mid], ALGO_CATALOG[mid:]]):
        x = MX if col_i == 0 else RIGHT_X
        y = MY_BODY
        for name, desc in chunk:
            _text(s, x, y, HALF_W, Inches(0.26), name, size=FONT.main, bold=True, color=BLACK)
            _text(s, x, y + Inches(0.24), HALF_W, Inches(0.38), desc, size=FONT.main - 1, color=MUTED)
            y += Inches(0.72)
    _code_footer(s, code)


def _app_table(prs: Presentation, code: str, title: str):
    s = _blank(prs)
    _heading(s, title)
    rows = _parse_report_cards()
    if not rows:
        _text(s, MX, Inches(1.5), CONTENT_W, Inches(1), "report_cards.md не найден", color=MUTED)
        _code_footer(s, code)
        return
    headers_ru = [
        "Алгоритм",
        "Покрытие static",
        "Покрытие dynamic",
        "Путь static, м",
        "Путь dynamic, м",
        "Наруш. у пешеходов (dynamic)",
    ]
    ncols = min(len(headers_ru), len(rows[0]))
    tbl = s.shapes.add_table(
        len(rows) + 1, ncols, MX, MY_BODY,
        CONTENT_W, Inches(0.38 * (len(rows) + 2)),
    ).table
    for j, h in enumerate(headers_ru[:ncols]):
        c = tbl.cell(0, j)
        c.text = h
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(FONT.table_header)
            p.font.bold = True
            p.font.color.rgb = BLACK
    for i, row in enumerate(rows, 1):
        for j in range(ncols):
            tbl.cell(i, j).text = row[j].strip("`") if j < len(row) else ""
            for p in tbl.cell(i, j).text_frame.paragraphs:
                p.font.size = Pt(FONT.table_body)
    _caption(s, "Медианы по seed · presentation_static + presentation_dynamic")
    _code_footer(s, code)


def _app_dual(
    prs: Presentation,
    code: str,
    title: str,
    left: Path | None,
    right: Path | None,
    cap: str = "",
    *,
    allow_gif: bool = True,
):
    s = _blank(prs)
    _heading(s, title)
    half = Inches(6.15)
    _img(s, left, MX, MY_BODY, HALF_W, Inches(5.55), allow_gif=allow_gif)
    _img(s, right, RIGHT_X, MY_BODY, HALF_W, Inches(5.55), allow_gif=allow_gif)
    _caption(s, cap or f"{_rel(left)}  ·  {_rel(right)}")
    _code_footer(s, code)


def _app_img(prs: Presentation, code: str, title: str, path: Path | None, cap: str = "", *, allow_gif: bool = False):
    s = _blank(prs)
    _heading(s, title)
    if not _img(s, path, MX, MY_BODY, CONTENT_W, Inches(5.55), allow_gif=allow_gif):
        _text(s, MX, Inches(3.2), CONTENT_W, Inches(0.5), f"Нет файла: {cap or _rel(path)}", color=MUTED)
    _caption(s, cap or _rel(path))
    _code_footer(s, code)


def _app_quad(prs: Presentation, code: str, title: str):
    s = _blank(prs)
    _heading(s, title)
    items = [
        ("frontier", A["flc_frontier_2d"]),
        ("voronoi", A["flc_voronoi_2d"]),
        ("darp", A["flc_darp_2d"]),
        ("stc", A["flc_stc_2d"]),
    ]
    pos = [(MX, MY_BODY), (RIGHT_X, MY_BODY), (MX, Inches(4.05)), (RIGHT_X, Inches(4.05))]
    labels = {
        "frontier": "baseline_frontier",
        "voronoi": "baseline_voronoi",
        "darp": "darp_boustro",
        "stc": "stc",
    }
    for (name, path), (l, t) in zip(items, pos):
        _img(s, path, l, t, HALF_W, Inches(2.65))
        _fig_label(s, l, t + Inches(2.68), HALF_W, labels.get(name, name))
    _code_footer(s, code)


def build_appendix_compact(prs: Presentation, *, embed_gif: bool):
    """Подвал: компактно, по темам. Коды приложений — кириллица (А, Б, В…)."""

    _app_img(prs, "А", "Приложение А · схема 5 слоёв", A["sch_classification"])

    _app_algo_catalog(prs, "Б", "Приложение Б · алгоритмы и краткое описание")

    _slide_table(
        prs, "Приложение Б1 · сложность и компромиссы",
        ["алгоритм", "время / шаг", "память", "плюс", "минус"],
        [
            ["baseline_random_walk", "O(1)",          "O(G)",        "минимум кода",     "низкое покрытие"],
            ["baseline_grid",        "O(G)",          "O(G)",        "детерминизм",      "без координации"],
            ["baseline_voronoi",     "O(N·G)",        "O(G)",        "зоны, меньше наруш.", "не адаптируется"],
            ["baseline_frontier",    "O(G+E) (BFS)",  "O(G)",        "быстрый рост покр.", "риск у людей"],
            ["boustrophedon",        "O(G)",          "O(G)",        "системный обход",  "только static"],
            ["stc",                  "O(G log G)",    "O(G)",        "полнота на сетке", "медленно в dynamic"],
            ["darp_boustro",         "O(I·N·G)",      "O(G)",        "равная нагрузка роботов", "длинный путь"],
            ["darp_stc",             "O(I·N·G log G)","O(G)",        "баланс + полнота", "наибольшая стоимость"],
            ["cbs_* (де-конфликт)",  "экспон. O(b^d)","O(N·G)",      "снимает коллизии", "только надстройка"],
            ["ppo_policy",           "O(W) inference","O(W)",        "RL-адаптация",     "обучение ~часы"],
            ["ml_guided / ml_goal",  "O(W)+O(G)",    "O(W+G)",      "адаптация + ограничитель","дороже классики"],
        ],
        col_widths=[2.0, 1.6, 0.9, 1.9, 1.9],
        caption="G — клетки сетки, N — роботы, E — рёбра, I — итерации, W — веса сети, b·d — ветвление и глубина CBS",
        notes="60 с. Главный справочный слайд по алгоритмам.",
    )

    # B2 · детали Frontier + Voronoi
    _slide_algo_detail(
        prs,
        "Приложение Б2 · Frontier и Voronoi",
        "Frontier",
        A["sch_frontier"],
        A["flc_frontier_2d"] or _ideal("baseline_frontier"),
        [
            "Для чего: разведка, неполная карта",
            "Как: граница visited / unvisited",
            "Числа: 87.9 с, наруш. у пешеходов 26",
        ],
        "Voronoi",
        A["sch_voronoi"],
        A["flc_voronoi_2d"] or _ideal("baseline_voronoi"),
        [
            "Для чего: параллельное покрытие зонами",
            "Как: разбиение между роботами",
            "Числа: 79.2 с, наруш. у пешеходов 32",
        ],
    )

    # B3 · STC и DARP
    _slide_algo_detail(
        prs,
        "Приложение Б3 · STC и DARP",
        "STC",
        A["sch_stc"],
        A["flc_stc_2d"] or _ideal("stc"),
        [
            "Для чего: полнота на static-сетке",
            "Как: остовное дерево",
            "Числа: 256 с, blocked 12 707",
        ],
        "DARP",
        A["sch_darp"],
        A["flc_darp_2d"] or _ideal("darp_boustro"),
        [
            "Для чего: баланс нагрузки",
            "Как: итеративное разбиение",
            "Числа: нагрузка равна (коэф. 0.01), путь 10 357 м",
        ],
    )

    # B4 · CBS и ML
    _slide_algo_detail(
        prs,
        "Приложение Б4 · CBS и ML-guided",
        "CBS",
        A["sch_cbs"],
        None,
        [
            "Для чего: де-конфликт роботов",
            "Как: поверх любого coverage",
            "Не строит покрытие сам",
        ],
        "ML + guardrail",
        A["sch_ml"],
        A["flc_ml_2d"] or A["ideal_guarded_2d"],
        [
            "Для чего: адаптация к сцене",
            "Как: SmallCNN + guardrail",
            "Режимы: чистый ML / мягкий / с ограничителем",
        ],
    )

    # C1 · Static сцена
    _slide_static_scene(prs, embed_gif)

    # C2 · Dynamic сцена
    _slide_dynamic_scene(prs, embed_gif)

    # D · Карты large_complex (4 алгоритма)
    _app_quad(prs, "Г", "Приложение Г · карты large_complex (покрытие)")

    _app_table(prs, "Д", "Приложение Д · сводка метрик (медианы по seed)")


def _slide_static_scene(prs: Presentation, embed_gif: bool):
    """Подвал · C1: static_A_long — суть сцены + GIF/2D + числа."""
    s = _blank(prs)
    _heading(s, "Приложение В1 · сцена static_A_long")
    _bullets(s, MX, MY_BODY, HALF_W - Inches(0.05), Inches(2.5), [
        "Цель: 0.90 coverage",
        "4 робота · без пешеходов",
        "Прямоугольные препятствия, известная карта",
        "Frontier: cov 0.90, путь 965 м, ped — нет",
        "Voronoi: cov 0.90, путь 977 м, дет.",
        "STC: cov 0.90, путь 1591 м",
        "DARP: cov 0.67 (порог не достигнут)",
    ], size=FONT.main - 1)
    img1 = A["stat_frontier_gif"] if embed_gif and A["stat_frontier_gif"] else A["stat_frontier_2d"]
    img2 = A["stat_voronoi_gif"] if embed_gif and A["stat_voronoi_gif"] else A["stat_voronoi_2d"]
    _img(s, img1, RIGHT_X, MY_BODY, HALF_W, Inches(2.65), allow_gif=embed_gif)
    _img(s, img2, RIGHT_X, MY_BODY + Inches(2.85), HALF_W, Inches(2.65), allow_gif=embed_gif)
    _fig_label(s, RIGHT_X, MY_BODY + Inches(2.62), HALF_W, "baseline_frontier")
    _fig_label(s, RIGHT_X, MY_BODY + Inches(2.62) + Inches(2.85), HALF_W, "baseline_voronoi")
    _code_footer(s, "В1")


def _slide_dynamic_scene(prs: Presentation, embed_gif: bool):
    """Подвал · C2: dynamic_B_long — суть сцены + GIF-сетка 5 алгоритмов + числа."""
    s = _blank(prs)
    _heading(s, "Приложение В2 · сцена dynamic_B_long")
    if embed_gif:
        items = [
            ("baseline_frontier", A["dyn_frontier_gif"]),
            ("baseline_voronoi", A["dyn_voronoi_gif"]),
            ("darp_boustro", A["dyn_darp_gif"]),
            ("stc", A["dyn_stc_gif"]),
            ("ml_guided", A["dyn_ml_gif"]),
        ]
        cols = 3
        cw = Inches(3.55)
        ch = Inches(2.25)
        g = Inches(0.18)
        for i, (label, path) in enumerate(items):
            c = i % cols
            r = i // cols
            l = MX + c * (cw + g)
            t = MY_BODY + r * (ch + Inches(0.45))
            _img(s, path, l, t, cw, ch - Inches(0.22), allow_gif=True)
            _fig_label(s, l, t + ch - Inches(0.24), cw, label)
        # Цифры справа от 2-го ряда
        _bullets(s, MX + 2 * (cw + g), MY_BODY + Inches(2.7), cw, Inches(2.2), [
            "Цель: 0.80, 4 робота, 6 пешеходов",
            "frontier: 0.80, наруш. 33",
            "voronoi: 0.80, наруш. 23",
            "ml_guided: 0.80 с ограничителем",
            "stc: блокировки в коридорах",
        ], size=FONT.main - 3)
    else:
        _img(s, A["dyn_frontier_png"], MX, MY_BODY, HALF_W, Inches(5.55))
        _bullets(s, RIGHT_X, MY_BODY, HALF_W, Inches(5.55), [
            "GIF: --embed-gif при сборке",
            "Цель: 0.80, 4 робота, 6 пешеходов",
            "Voronoi: наруш. 23 vs frontier 33",
            "ml_guided с ограничителем — мягче по покрытию",
        ])
    _code_footer(s, "В2")


def build_appendix_full(prs: Presentation, *, embed_gif: bool):
    """Расширенный подвал: compact + доп. графики."""
    build_appendix_compact(prs, embed_gif=embed_gif)
    _app_img(prs, "C+", "coverage vs time", A["flc_frontier_time"], allow_gif=False)
    _app_img(prs, "C++", "distance vs coverage", A["flc_frontier_dist"], allow_gif=False)


def build(
    font: FontProfile = FONT_MAIN14,
    out: Path = OUT,
    *,
    compact: bool = True,
    embed_gif: bool = False,
    full_appendix: bool = False,
) -> Path:
    global FONT
    FONT = font

    if embed_gif and not compact:
        print("Note: GIF embedding increases file size (tens of MB).")

    missing = [k for k, v in A.items() if v is None and k != "report_cards"]
    if missing:
        print("Warning missing assets:", ", ".join(missing))

    if Image is None:
        print("Tip: pip install Pillow for smaller JPEG embeds")

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    build_main(prs, embed_gif=embed_gif)
    if full_appendix:
        build_appendix_full(prs, embed_gif=embed_gif)
    else:
        build_appendix_compact(prs, embed_gif=embed_gif)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    n = len(prs.slides)
    mb = out.stat().st_size / (1024 * 1024)
    mode = "compact+gif" if embed_gif and not full_appendix else ("full" if full_appendix else "compact")
    gif_n = sum(1 for v in A.values() if v and str(v).lower().endswith(".gif"))
    print(f"Saved: {out}")
    print(f"Slides: {n} · Size: {mb:.2f} MB · mode={mode} · gif_assets={gif_n}")
    return out


def main():
    p = argparse.ArgumentParser(description="Build defense PPTX")
    p.add_argument("--default-font", action="store_true", help="Larger fonts (17/11)")
    p.add_argument("-o", "--output", type=Path, default=OUT)
    p.add_argument("--embed-gif", action="store_true", help="Embed GIF in slides (large file)")
    p.add_argument("--full-appendix", action="store_true", help="Extra appendix charts (+GIF slide if --embed-gif)")
    args = p.parse_args()
    font = FONT_DEFAULT if args.default_font else FONT_MAIN14
    build(
        font=font,
        out=args.output,
        compact=not args.full_appendix,
        embed_gif=args.embed_gif,
        full_appendix=args.full_appendix,
    )


if __name__ == "__main__":
    main()
